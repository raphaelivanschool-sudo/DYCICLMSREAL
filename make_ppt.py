# -*- coding: utf-8 -*-
"""Generate the DYCI Computer Laboratory Management System presentation."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---- Theme ----------------------------------------------------------------
NAVY   = RGBColor(0x0F, 0x2A, 0x4A)   # primary dark
BLUE   = RGBColor(0x2C, 0x7B, 0xE5)   # accent
LIGHT  = RGBColor(0xF2, 0xF6, 0xFC)   # light panel
GREY   = RGBColor(0x5B, 0x6B, 0x7B)   # body text
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GOLD   = RGBColor(0xF5, 0xA6, 0x23)   # secondary accent

FONT = "Segoe UI"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_slide():
    return prs.slides.add_slide(BLANK)


def rect(slide, x, y, w, h, color, line=None):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp


def textbox(slide, x, y, w, h, lines, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        text, size, color, bold = ln[0], ln[1], ln[2], ln[3]
        align = ln[4] if len(ln) > 4 else PP_ALIGN.LEFT
        space = ln[5] if len(ln) > 5 else 4
        p.alignment = align
        p.space_after = Pt(space)
        p.space_before = Pt(0)
        r = p.add_run()
        r.text = text
        r.font.name = FONT
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.bold = bold
    return tb


def header(slide, kicker, title):
    """Standard content-slide header with side accent bar."""
    rect(slide, 0, 0, SW, Inches(1.35), NAVY)
    rect(slide, 0, 0, Inches(0.18), Inches(1.35), GOLD)
    textbox(slide, Inches(0.55), Inches(0.18), Inches(12), Inches(1.0),
            [(kicker, 12, RGBColor(0x9D, 0xC2, 0xF0), True, PP_ALIGN.LEFT, 2),
             (title, 28, WHITE, True, PP_ALIGN.LEFT, 0)])


def bullet_card(slide, x, y, w, h, title, items, accent=BLUE):
    rect(slide, x, y, w, h, LIGHT)
    rect(slide, x, y, w, Inches(0.07), accent)
    tb = slide.shapes.add_textbox(x + Inches(0.22), y + Inches(0.2),
                                  w - Inches(0.4), h - Inches(0.35))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.space_after = Pt(6)
    r = p.add_run(); r.text = title
    r.font.name = FONT; r.font.size = Pt(15); r.font.bold = True; r.font.color.rgb = NAVY
    for it in items:
        p = tf.add_paragraph()
        p.space_after = Pt(3)
        r = p.add_run(); r.text = "•  " + it
        r.font.name = FONT; r.font.size = Pt(11); r.font.color.rgb = GREY
    return tb


# ===========================================================================
# Slide 1 — Title
# ===========================================================================
s = add_slide()
rect(s, 0, 0, SW, SH, NAVY)
rect(s, 0, Inches(5.35), SW, Inches(0.12), GOLD)
# accent block
rect(s, Inches(0.9), Inches(2.05), Inches(0.9), Inches(0.12), BLUE)
textbox(s, Inches(0.9), Inches(2.25), Inches(11.5), Inches(3.0),
        [("DYCI Computer Laboratory", 44, WHITE, True, PP_ALIGN.LEFT, 2),
         ("Management System", 44, WHITE, True, PP_ALIGN.LEFT, 10),
         ("A centralized platform to monitor, control, and manage "
          "computer laboratories in real time.", 17,
          RGBColor(0xB9, 0xD0, 0xEE), False, PP_ALIGN.LEFT, 0)])
textbox(s, Inches(0.9), Inches(5.7), Inches(11.5), Inches(1.0),
        [("Capstone / Project Presentation", 13, GOLD, True, PP_ALIGN.LEFT, 2),
         ("Dr. Yanga's Colleges, Inc.  •  Computer Laboratory", 12,
          RGBColor(0x9D, 0xB6, 0xD4), False, PP_ALIGN.LEFT, 0)])

# ===========================================================================
# Slide 2 — Project Overview
# ===========================================================================
s = add_slide()
header(s, "INTRODUCTION", "Project Overview")
textbox(s, Inches(0.55), Inches(1.7), Inches(12.2), Inches(1.4),
        [("The DYCI Computer Laboratory Management System is a web-based platform that lets "
          "administrators and instructors monitor and manage every PC in a computer lab from a "
          "single dashboard — turning manual, hands-on lab supervision into one centralized control point.",
          15, GREY, False, PP_ALIGN.LEFT, 0)])

cards = [
    ("Centralized Control", ["One dashboard for all lab PCs",
                             "Lock, shut down & restrict access",
                             "Live online/offline status"], BLUE),
    ("Real-Time & Connected", ["Live screen monitoring",
                               "Instant screen projection",
                               "Built on Socket.IO updates"], GOLD),
    ("Academics Built-In", ["Grading & gradebook",
                            "Messaging & support tickets",
                            "Usage & activity reports"], BLUE),
]
cw = Inches(3.95); gap = Inches(0.25); x0 = Inches(0.55); cy = Inches(3.25); ch = Inches(3.0)
for i, (t, items, ac) in enumerate(cards):
    bullet_card(s, x0 + i * (cw + gap), cy, cw, ch, t, items, ac)

# ===========================================================================
# Slide 3 — How it works / Architecture
# ===========================================================================
s = add_slide()
header(s, "AT A GLANCE", "How the System Works")
textbox(s, Inches(0.55), Inches(1.65), Inches(12.2), Inches(0.7),
        [("Two connected sides keep the whole lab in sync over the local network.",
          15, GREY, False, PP_ALIGN.LEFT, 0)])

# Host box
def flow_box(slide, x, y, w, h, title, sub, color):
    rect(slide, x, y, w, h, color)
    textbox(slide, x + Inches(0.2), y + Inches(0.18), w - Inches(0.4), h - Inches(0.3),
            [(title, 16, WHITE, True, PP_ALIGN.LEFT, 3),
             (sub, 11, RGBColor(0xDD, 0xE9, 0xF7), False, PP_ALIGN.LEFT, 0)],
            anchor=MSO_ANCHOR.MIDDLE)

fy = Inches(2.85); fh = Inches(1.5)
flow_box(s, Inches(0.55), fy, Inches(3.7), fh, "Host Dashboard",
         "Admin / Instructor controls the lab from the browser", NAVY)
flow_box(s, Inches(4.85), fy, Inches(3.7), fh, "Host Server",
         "Relays commands & screen frames securely", BLUE)
flow_box(s, Inches(9.15), fy, Inches(3.6), fh, "Lab PC Agents",
         "Each PC reports status & runs locked overlay", NAVY)
# arrows
for ax in (Inches(4.32), Inches(8.62)):
    ar = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, ax, fy + Inches(0.55), Inches(0.45), Inches(0.4))
    ar.fill.solid(); ar.fill.fore_color.rgb = GOLD; ar.line.fill.background(); ar.shadow.inherit = False

textbox(s, Inches(0.55), Inches(5.0), Inches(12.2), Inches(1.8),
        [("Key capabilities", 14, NAVY, True, PP_ALIGN.LEFT, 6),
         ("•  Live monitoring of every lab PC — online status, screens & activity",
          12, GREY, False, PP_ALIGN.LEFT, 4),
         ("•  Locked Demo Mode — project the instructor's screen fullscreen and input-locked on guests",
          12, GREY, False, PP_ALIGN.LEFT, 4),
         ("•  Remote actions — lock, shut down, network scan & website restriction",
          12, GREY, False, PP_ALIGN.LEFT, 0)])

# ===========================================================================
# Slide 4 — Target Users
# ===========================================================================
s = add_slide()
header(s, "WHO IT'S FOR", "Target Users")
roles = [
    ("Administrator", "Manages the whole system",
     ["Users, labs & computers", "Network & security control",
      "Reports, logs & inventory"], BLUE),
    ("Instructor", "Runs the classroom",
     ["Monitor student screens", "Lock / project / control PCs",
      "Grade students & approve tickets"], GOLD),
    ("Student", "Uses the lab",
     ["View session dashboard", "Message & request support",
      "Check personal grades"], BLUE),
]
cw = Inches(3.95); gap = Inches(0.25); x0 = Inches(0.55); cy = Inches(2.0); ch = Inches(4.4)
for i, (name, sub, items, ac) in enumerate(roles):
    x = x0 + i * (cw + gap)
    rect(s, x, cy, cw, ch, LIGHT)
    rect(s, x, cy, cw, Inches(0.9), NAVY)
    textbox(s, x + Inches(0.25), cy + Inches(0.16), cw - Inches(0.5), Inches(0.7),
            [(name, 18, WHITE, True, PP_ALIGN.LEFT, 1),
             (sub, 11, RGBColor(0xB9, 0xD0, 0xEE), False, PP_ALIGN.LEFT, 0)])
    rect(s, x, cy + Inches(0.9), cw, Inches(0.06), ac)
    tb = s.shapes.add_textbox(x + Inches(0.25), cy + Inches(1.2), cw - Inches(0.5), ch - Inches(1.4))
    tf = tb.text_frame; tf.word_wrap = True
    for j, it in enumerate(items):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.space_after = Pt(10)
        r = p.add_run(); r.text = "•  " + it
        r.font.name = FONT; r.font.size = Pt(13); r.font.color.rgb = GREY

# ===========================================================================
# Slide 5 — Admin modules
# ===========================================================================
s = add_slide()
header(s, "MODULES & FEATURES — ADMINISTRATOR", "What It Does for Admins")
admin = [
    ("User & Access Management", ["Manage admin, instructor & student accounts",
                                  "Role-based access control"]),
    ("Lab & Computer Management", ["Register laboratories & computers",
                                   "Hardware inventory tracking"]),
    ("Scheduling", ["Assign labs, classes & time slots",
                    "Avoid booking conflicts"]),
    ("Network & Security", ["Scan the network & restrict websites",
                            "Security settings & policies"]),
    ("Monitoring & Projection", ["Live PC status & screen projection",
                                 "Lock / shut down any PC"]),
    ("Reports & System Logs", ["Usage & activity reports",
                               "Full audit trail of actions"]),
]
cw = Inches(3.95); ch = Inches(2.35); gx = Inches(0.25); gy = Inches(0.25)
x0 = Inches(0.55); y0 = Inches(1.65)
for i, (t, items) in enumerate(admin):
    r, c = divmod(i, 3)
    bullet_card(s, x0 + c * (cw + gx), y0 + r * (ch + gy), cw, ch, t, items, BLUE)

# ===========================================================================
# Slide 6 — Instructor modules
# ===========================================================================
s = add_slide()
header(s, "MODULES & FEATURES — INSTRUCTOR", "What It Does for Instructors")
inst = [
    ("Classroom Dashboard", ["See the whole class at a glance",
                             "Live student status"]),
    ("Student Screen Monitoring", ["View student screens in real time",
                                   "Spot off-task activity"]),
    ("Control Actions", ["Lock, shut down & project screen",
                         "Keep the class focused"]),
    ("Messaging", ["Send announcements to students",
                   "One-to-one chat support"]),
    ("Ticket Approval", ["Review student support requests",
                         "Approve or resolve quickly"]),
    ("Grading & Reports", ["Gradebook, grade scales & subjects",
                           "Class performance reports"]),
]
for i, (t, items) in enumerate(inst):
    r, c = divmod(i, 3)
    bullet_card(s, x0 + c * (cw + gx), y0 + r * (ch + gy), cw, ch, t, items, GOLD)

# ===========================================================================
# Slide 7 — Student modules
# ===========================================================================
s = add_slide()
header(s, "MODULES & FEATURES — STUDENT", "What It Does for Students")
stud = [
    ("Session Dashboard", ["See current class & session info",
                           "Stay aware of lab activity"]),
    ("Support Tickets", ["Report PC or software issues",
                         "Track request status"]),
    ("Messaging", ["Chat with the instructor",
                   "Receive announcements"]),
    ("My Grades", ["View personal grades anytime",
                   "Transparent grade breakdown"]),
]
cw2 = Inches(5.95); x02 = Inches(0.55)
for i, (t, items) in enumerate(stud):
    r, c = divmod(i, 2)
    bullet_card(s, x02 + c * (cw2 + Inches(0.35)), y0 + r * (ch + gy + Inches(0.15)),
                cw2, Inches(2.5), t, items, BLUE)

# ===========================================================================
# Slide 8 — Benefits / Closing
# ===========================================================================
s = add_slide()
rect(s, 0, 0, SW, SH, NAVY)
rect(s, Inches(0.9), Inches(0.85), Inches(0.9), Inches(0.12), GOLD)
textbox(s, Inches(0.9), Inches(1.05), Inches(11.5), Inches(1.0),
        [("KEY BENEFITS", 13, GOLD, True, PP_ALIGN.LEFT, 3),
         ("Why It Matters", 34, WHITE, True, PP_ALIGN.LEFT, 0)])

bens = [
    ("Saves time", "Manage every PC from one screen — no walking station to station."),
    ("Better focus", "Lock screens & project demos to keep students on task."),
    ("Faster support", "Tickets & messaging resolve lab issues quickly."),
    ("Clear records", "Grades, reports & logs keep everything accountable."),
]
bw = Inches(5.7); bh = Inches(1.6); bx = Inches(0.9); by = Inches(2.7)
for i, (t, d) in enumerate(bens):
    r, c = divmod(i, 2)
    x = bx + c * (bw + Inches(0.4)); y = by + r * (bh + Inches(0.35))
    rect(s, x, y, Inches(0.07), bh, GOLD)
    textbox(s, x + Inches(0.3), y + Inches(0.1), bw - Inches(0.4), bh - Inches(0.2),
            [(t, 18, WHITE, True, PP_ALIGN.LEFT, 4),
             (d, 13, RGBColor(0xB9, 0xD0, 0xEE), False, PP_ALIGN.LEFT, 0)],
            anchor=MSO_ANCHOR.TOP)

textbox(s, Inches(0.9), Inches(6.65), Inches(11.5), Inches(0.6),
        [("DYCI Computer Laboratory Management System  —  smarter labs, simpler management.",
          13, RGBColor(0x9D, 0xB6, 0xD4), True, PP_ALIGN.LEFT, 0)])

out = r"c:\Users\acer\Documents\GitHub\DYCICLMS45\DYCI_CLMS_Presentation.pptx"
prs.save(out)
print("Saved:", out)
print("Slides:", len(prs.slides._sldIdLst))
